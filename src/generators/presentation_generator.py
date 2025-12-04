# src/generators/presentation_generator.py
from typing import List, Dict, Any
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
from pptx import Presentation
from pptx.util import Inches
import json

class PresentationGenerator:
    def __init__(self, settings):
        self.settings = settings
        self.llm = ChatOpenAI(
            model_name=settings.summary_model,
            temperature=0.3
        )
    
    async def generate_presentation(self, research_data: Dict[str, Any]) -> str:
        """Generate a PowerPoint presentation from research data"""
        
        # Generate slide content structure
        slide_structure = await self._create_slide_structure(research_data)
        
        # Create PowerPoint presentation
        presentation = self._create_pptx_presentation(slide_structure)
        
        # Save presentation
        filename = f"research_presentation_{research_data['research_plan'].topic.replace(' ', '_')}.pptx"
        presentation.save(filename)
        
        return filename
    
    async def _create_slide_structure(self, research_data: Dict[str, Any]) -> List[Dict]:
        """Create structured content for each slide"""
        
        system_prompt = """You are an expert at creating executive presentations. 
        Create a slide-by-slide structure for a technical executive summary."""
        
        user_prompt = f"""
        Research Topic: {research_data['research_plan'].topic}
        
        Synthesized Content:
        {json.dumps(research_data['synthesized_content'], indent=2)}
        
        Create {self.settings.max_slides} slides maximum with this structure:
        - Title Slide
        - Executive Summary
        - Key Findings
        - Technical Overview
        - Current State Analysis
        - Future Outlook
        - Challenges & Opportunities
        - Recommendations
        - Conclusion
        
        For each slide, provide:
        - Title
        - 3-5 bullet points
        - Suggested visualizations (if any)
        """
        
        response = await self.llm.agenerate([[SystemMessage(content=system_prompt), 
                                            HumanMessage(content=user_prompt)]])
        
        return self._parse_slide_structure(response.generations[0][0].text)
    
    def _create_pptx_presentation(self, slide_structure: List[Dict]) -> Presentation:
        """Create actual PowerPoint file"""
        
        prs = Presentation()
        
        for slide_data in slide_structure:
            # Choose slide layout based on content
            if slide_data['title'].lower() == 'title slide':
                slide_layout = prs.slide_layouts[0]  # Title slide
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            # Add content
            if hasattr(slide.shapes, 'placeholders') and len(slide.shapes.placeholders) > 1:
                content_shape = slide.shapes.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.text = ""
                
                for bullet in slide_data['bullet_points']:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
        
        return prs